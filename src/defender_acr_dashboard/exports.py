from __future__ import annotations

from pathlib import Path
from typing import Any

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from .config import DEFENDER_SERVICE, TOTAL_SERVICE
from .dashboard_model import build_dashboard_model


DEFAULT_EXPORT_THRESHOLD = 8.0

BLUE = "0078D4"
NAVY = "0F3A5F"
DARK = "201F1E"
TEXT = "323130"
MUTED = "605E5C"
LIGHT_BG = "F5F7FB"
BORDER = "D9E2EC"
RED = "D13438"
ORANGE = "FF8C00"
GREEN = "107C10"
WHITE = "FFFFFF"


def create_powerpoint(
    records: pd.DataFrame,
    model_or_metrics: dict[str, Any] | pd.DataFrame,
    source_name: str,
    output_path: Path,
    defender_share_threshold: float = DEFAULT_EXPORT_THRESHOLD,
) -> Path:
    model = model_or_metrics if isinstance(model_or_metrics, dict) else build_dashboard_model(records)
    rows = [row for row in model["opportunity"] if row["opportunity"] != "Too small"]
    actions = _action_rows(model, defender_share_threshold)
    portfolio = _portfolio_summary(model, defender_share_threshold)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.title = "Defender for Cloud ACR opportunities"

    _title_slide(prs, source_name, model, portfolio, defender_share_threshold)
    _portfolio_slide(prs, model, portfolio, actions, defender_share_threshold)
    _action_queue_slide(prs, actions, defender_share_threshold)
    _opportunity_slide(prs, rows)

    for row in actions[:3]:
        _customer_slide(prs, model, row, defender_share_threshold)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return output_path


def _title_slide(
    prs: Presentation,
    source_name: str,
    model: dict[str, Any],
    portfolio: dict[str, float],
    threshold: float,
) -> None:
    slide = _blank_slide(prs, NAVY)
    _add_text(slide, "Defender for Cloud ACR opportunities", 0.55, 0.55, 8.6, 0.65, 32, WHITE, bold=True)
    _add_text(slide, "Executive opportunity readout", 0.58, 1.25, 5.0, 0.35, 15, "C7E0F4")
    _add_text(slide, f"Source workbook: {source_name}", 0.58, 6.75, 7.5, 0.25, 9, "C7E0F4")
    _add_text(slide, f"Latest month: {model.get('last_full_month', '-')}", 9.75, 0.65, 2.6, 0.3, 11, "C7E0F4", align=PP_ALIGN.RIGHT)

    _add_stat_card(slide, "Annualized DfC opportunity", _money(portfolio["annual_opportunity"]), 0.65, 2.25, 3.7, 1.5, RED)
    _add_stat_card(slide, "Customers below threshold", f"{int(portfolio['below_threshold'])}", 4.75, 2.25, 3.7, 1.5, ORANGE)
    _add_stat_card(slide, "Portfolio DfC share", _pct(portfolio["monthly_dfc_share"]), 8.85, 2.25, 3.7, 1.5, GREEN)

    narrative = (
        f"At a {threshold:.0f}% Defender share threshold, the visible portfolio has "
        f"{int(portfolio['below_threshold'])} customers below target and an estimated "
        f"{_money(portfolio['annual_opportunity'])} annualized DfC ACR run-rate gap."
    )
    _add_text(slide, narrative, 0.75, 4.35, 11.7, 0.8, 20, WHITE, bold=True)
    _add_text(
        slide,
        "This deck uses latest monthly ACR for run-rate opportunity sizing and FY-to-date ACR for total-account context.",
        0.75,
        5.25,
        11.2,
        0.35,
        12,
        "C7E0F4",
    )


def _portfolio_slide(
    prs: Presentation,
    model: dict[str, Any],
    portfolio: dict[str, float],
    actions: list[dict[str, Any]],
    threshold: float,
) -> None:
    slide = _blank_slide(prs)
    _slide_title(slide, "Portfolio snapshot")

    _add_stat_card(slide, "FYTD Total ACR", _money(portfolio["fytd_total"]), 0.55, 1.2, 2.9, 1.0, BLUE)
    _add_stat_card(slide, "FYTD DfC ACR", _money(portfolio["fytd_dfc"]), 3.7, 1.2, 2.9, 1.0, GREEN)
    _add_stat_card(slide, "Monthly Total ACR", _money(portfolio["monthly_total"]), 6.85, 1.2, 2.9, 1.0, BLUE)
    _add_stat_card(slide, "Monthly DfC ACR", _money(portfolio["monthly_dfc"]), 10.0, 1.2, 2.9, 1.0, GREEN)

    _add_text(slide, "Executive narrative", 0.65, 2.65, 3.5, 0.3, 16, DARK, bold=True)
    bullets = [
        f"Default export threshold: {threshold:.0f}% Defender share of latest monthly Total ACR.",
        f"Annualized opportunity is {_money(portfolio['annual_opportunity'])}, calculated from monthly gap to threshold times 12.",
        f"Top action account: {actions[0]['customer'] if actions else '-'} with {_money(actions[0]['annual_opportunity']) if actions else '$0'} annualized opportunity.",
        "Opportunity priority keeps the dashboard logic: low Defender share, meaningful Azure footprint, and growth gap.",
    ]
    _add_bullets(slide, bullets, 0.75, 3.05, 5.45, 2.2)

    labels = model.get("month_labels", [])
    total_series = model.get("product_monthly", {}).get(TOTAL_SERVICE, [])
    dfc_series = model.get("dfc_total_monthly", [])
    _add_line_chart(
        slide,
        "Monthly ACR trend",
        labels,
        [("Total", total_series, BLUE), ("Defender", dfc_series, GREEN)],
        6.75,
        2.65,
        5.75,
        2.5,
    )
    _add_text(
        slide,
        "Note: FYTD values sum months in the latest fiscal year; monthly values are the latest available month.",
        0.75,
        6.65,
        11.5,
        0.25,
        9,
        MUTED,
    )


def _action_queue_slide(prs: Presentation, actions: list[dict[str, Any]], threshold: float) -> None:
    slide = _blank_slide(prs)
    _slide_title(slide, "Sales action queue")
    _add_text(
        slide,
        f"Prioritized at {threshold:.0f}% Defender share. Annualized opportunity is a run-rate estimate, not forecast or pipeline.",
        0.65,
        0.95,
        11.4,
        0.3,
        11,
        MUTED,
    )
    headers = ["#", "Customer", "Priority", "FYTD Total", "Monthly DfC %", "Annual opp.", "Recommended action"]
    data = [
        [
            str(index),
            _short(row["customer"], 28),
            row["opportunity"],
            _money(row["total_fytd"]),
            _pct(row["dfc_ratio"] / 100),
            _money(row["annual_opportunity"]),
            _short(row["recommended_action"], 34),
        ]
        for index, row in enumerate(actions[:9], start=1)
    ]
    _add_table(
        slide,
        headers,
        data,
        0.45,
        1.45,
        12.45,
        5.45,
        [0.35, 2.45, 0.8, 1.35, 1.05, 1.25, 5.2],
        [2, 28, 8, 12, 10, 12, 34],
    )


def _opportunity_slide(prs: Presentation, rows: list[dict[str, Any]]) -> None:
    slide = _blank_slide(prs)
    _slide_title(slide, "Opportunity matrix details")
    _add_text(
        slide,
        "Both FYTD and latest monthly ACR are shown to avoid mixing cumulative account context with run-rate opportunity signals.",
        0.65,
        0.95,
        11.8,
        0.3,
        11,
        MUTED,
    )
    ordered = sorted(
        rows,
        key=lambda row: (
            0 if row["opportunity"] == "High" else 1 if row["opportunity"] == "Medium" else 2,
            -(row.get("growth_gap") or 0),
        ),
    )
    headers = ["Customer", "Priority", "FYTD Total", "Monthly Total", "FYTD DfC", "Monthly DfC", "DfC %", "Signal"]
    data = [
        [
            _short(row["customer"], 26),
            row["opportunity"],
            _money(row.get("total_fytd", 0)),
            _money(row.get("total_monthly_current", row.get("total_current", 0))),
            _money(row.get("dfc_fytd", 0)),
            _money(row.get("dfc_monthly_current", row.get("dfc_current", 0))),
            _pct(row.get("dfc_ratio", 0) / 100),
            _short(row.get("notes", "-"), 44),
        ]
        for row in ordered[:10]
    ]
    _add_table(
        slide,
        headers,
        data,
        0.35,
        1.35,
        12.65,
        5.8,
        [2.2, 0.75, 1.25, 1.25, 1.15, 1.15, 0.65, 4.25],
        [26, 8, 12, 12, 12, 12, 6, 44],
    )


def _customer_slide(prs: Presentation, model: dict[str, Any], row: dict[str, Any], threshold: float) -> None:
    slide = _blank_slide(prs)
    customer = row["customer"]
    data = model["customer_data"].get(customer, {})
    labels = model.get("month_labels", [])

    _slide_title(slide, customer)
    _add_text(slide, row["recommended_action"], 0.65, 0.95, 6.0, 0.35, 14, RED, bold=True)

    _add_stat_card(slide, "FYTD Total ACR", _money(row.get("total_fytd", 0)), 0.55, 1.45, 2.5, 0.9, BLUE)
    _add_stat_card(slide, "Monthly Total ACR", _money(row.get("total_monthly_current", 0)), 3.25, 1.45, 2.5, 0.9, BLUE)
    _add_stat_card(slide, "Monthly DfC share", _pct(row.get("dfc_ratio", 0) / 100), 5.95, 1.45, 2.5, 0.9, ORANGE)
    _add_stat_card(slide, "Annualized opp.", _money(row.get("annual_opportunity", 0)), 8.65, 1.45, 2.5, 0.9, RED)

    _add_line_chart(
        slide,
        "Customer monthly ACR trend",
        labels,
        [
            ("Total", data.get("total_series", []), BLUE),
            ("Defender", data.get("dfc_series", []), GREEN),
        ],
        0.65,
        2.75,
        5.8,
        2.25,
    )
    share_series = _share_series(data.get("dfc_series", []), data.get("total_series", []))
    _add_line_chart(
        slide,
        "Defender share trend",
        labels,
        [("DfC share", share_series, ORANGE)],
        6.85,
        2.75,
        5.8,
        2.25,
        percent_axis=True,
    )

    products = [
        product
        for product in data.get("products", [])
        if product.get("product") != DEFENDER_SERVICE and product.get("current", 0) > 0
    ][:5]
    bullets = [row["conversation_angle"], row["action_reason"]]
    bullets.extend(
        f"{product['product']}: {_money(product['current'])} monthly ACR, 3M {_signed_pct(product.get('three_m'))}"
        for product in products[:2]
    )
    _add_text(slide, "Recommended follow-up", 0.65, 5.45, 2.8, 0.3, 15, DARK, bold=True)
    _add_bullets(slide, bullets, 0.75, 5.85, 11.3, 1.25, font_size=9)
    _add_text(slide, f"Threshold basis: {threshold:.0f}% of latest monthly Total ACR.", 9.0, 0.95, 3.2, 0.25, 9, MUTED, align=PP_ALIGN.RIGHT)


def _portfolio_summary(model: dict[str, Any], threshold: float) -> dict[str, float]:
    rows = [row for row in model["opportunity"] if row["opportunity"] != "Too small"]
    fytd_total = sum(row.get("total_fytd", 0) for row in rows)
    fytd_dfc = sum(row.get("dfc_fytd", 0) for row in rows)
    monthly_total = sum(row.get("total_monthly_current", row.get("total_current", 0)) for row in rows)
    monthly_dfc = sum(row.get("dfc_monthly_current", row.get("dfc_current", 0)) for row in rows)
    monthly_gap = sum(_monthly_gap(row, threshold) for row in rows)
    return {
        "fytd_total": fytd_total,
        "fytd_dfc": fytd_dfc,
        "monthly_total": monthly_total,
        "monthly_dfc": monthly_dfc,
        "monthly_dfc_share": monthly_dfc / monthly_total if monthly_total else 0,
        "monthly_gap": monthly_gap,
        "annual_opportunity": monthly_gap * 12,
        "below_threshold": sum(1 for row in rows if row.get("dfc_ratio", 0) < threshold),
    }


def _action_rows(model: dict[str, Any], threshold: float) -> list[dict[str, Any]]:
    priority_order = {"High": 0, "Medium": 1, "Low": 2, "Too small": 3}
    rows = []
    for row in model["opportunity"]:
        if row["opportunity"] == "Too small":
            continue
        monthly_gap = _monthly_gap(row, threshold)
        action, angle, reason = _action_text(model, row, threshold)
        rows.append(
            {
                **row,
                "monthly_gap": monthly_gap,
                "annual_opportunity": monthly_gap * 12,
                "recommended_action": action,
                "conversation_angle": angle,
                "action_reason": reason,
            }
        )
    return sorted(
        rows,
        key=lambda row: (
            row.get("dfc_ratio", 0) >= threshold,
            priority_order[row["opportunity"]],
            -row["annual_opportunity"],
            -(row.get("growth_gap") or 0),
        ),
    )


def _monthly_gap(row: dict[str, Any], threshold: float) -> float:
    total = row.get("total_monthly_current", row.get("total_current", 0)) or 0
    dfc = row.get("dfc_monthly_current", row.get("dfc_current", 0)) or 0
    return max(0.0, total * (threshold / 100) - dfc)


def _action_text(model: dict[str, Any], row: dict[str, Any], threshold: float) -> tuple[str, str, str]:
    customer_data = model.get("customer_data", {}).get(row["customer"], {})
    products = customer_data.get("products", [])
    workload = next(
        (
            item["product"]
            for item in products
            if item.get("product") != DEFENDER_SERVICE and item.get("current", 0) > 0
        ),
        "core Azure workloads",
    )
    below_threshold = row.get("dfc_ratio", 0) < threshold
    if row.get("dfc_monthly_current", row.get("dfc_current", 0)) < 30 and row.get("total_monthly_current", 0) > 3000:
        return (
            "Start DfC attach discovery",
            f"Open with current {workload} usage and validate Defender for Cloud coverage.",
            "Little or no Defender for Cloud ACR against a meaningful Azure footprint.",
        )
    if below_threshold and (row.get("growth_gap") or 0) > 0:
        return (
            "Prioritize attach expansion",
            f"Lead with {workload} growth and the Defender share gap to threshold.",
            "Azure footprint is growing faster than Defender for Cloud attach.",
        )
    if below_threshold:
        return (
            "Expand Defender coverage",
            f"Review Defender for Cloud coverage across {workload} and adjacent services.",
            "Defender for Cloud share is below the selected threshold.",
        )
    return (
        "Monitor Defender attach",
        f"Confirm Defender for Cloud coverage keeps pace with {workload}.",
        row.get("notes") if row.get("notes") and row.get("notes") != "-" else "No urgent attach gap under the selected threshold.",
    )


def _blank_slide(prs: Presentation, background: str = WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, 13.333, 7.5, background, background)
    return slide


def _slide_title(slide, title: str) -> None:
    _add_text(slide, title, 0.55, 0.35, 9.6, 0.45, 24, DARK, bold=True)
    _add_rect(slide, 0.55, 0.92, 1.1, 0.06, BLUE, BLUE)


def _add_stat_card(slide, label: str, value: str, x: float, y: float, w: float, h: float, color: str) -> None:
    _add_rect(slide, x, y, w, h, WHITE, BORDER)
    _add_rect(slide, x, y, 0.08, h, color, color)
    _add_text(slide, label.upper(), x + 0.2, y + 0.14, w - 0.35, 0.22, 8, MUTED, bold=True)
    _add_text(slide, value, x + 0.2, y + 0.44, w - 0.35, 0.42, 20, color, bold=True)


def _add_table(
    slide,
    headers: list[str],
    rows: list[list[str]],
    x: float,
    y: float,
    w: float,
    h: float,
    col_widths: list[float],
    col_limits: list[int] | None = None,
) -> None:
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    for index, width in enumerate(col_widths):
        table.columns[index].width = Inches(width)
    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = header
        _style_cell(cell, fill=NAVY, color=WHITE, bold=True, font_size=8)
    for row_index, row in enumerate(rows, start=1):
        for col, value in enumerate(row):
            cell = table.cell(row_index, col)
            limit = col_limits[col] if col_limits and col < len(col_limits) else 48
            cell.text = _short(str(value), limit)
            fill = "F8FAFC" if row_index % 2 == 0 else WHITE
            color = TEXT
            bold = False
            if col < len(row) and value in {"High", "Medium", "Low"}:
                color = RED if value == "High" else ORANGE if value == "Medium" else GREEN
                bold = True
            _style_cell(cell, fill=fill, color=color, bold=bold, font_size=7)


def _style_cell(cell, *, fill: str, color: str, bold: bool = False, font_size: int = 8) -> None:
    cell.fill.solid()
    cell.fill.fore_color.rgb = _rgb(fill)
    cell.margin_left = Inches(0.05)
    cell.margin_right = Inches(0.05)
    text_frame = cell.text_frame
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = _rgb(color)
        paragraph.font.bold = bold


def _add_line_chart(
    slide,
    title: str,
    labels: list[str],
    series: list[tuple[str, list[float], str]],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    percent_axis: bool = False,
) -> None:
    _add_rect(slide, x, y, w, h, WHITE, BORDER)
    _add_text(slide, title, x + 0.18, y + 0.12, w - 0.3, 0.25, 12, DARK, bold=True)
    plot_x, plot_y, plot_w, plot_h = x + 0.35, y + 0.55, w - 0.7, h - 0.95
    all_values = [value for _, values, _ in series for value in values if value is not None]
    if not all_values:
        _add_text(slide, "No trend data", plot_x, plot_y + 0.5, plot_w, 0.3, 10, MUTED, align=PP_ALIGN.CENTER)
        return
    min_value = 0 if min(all_values) >= 0 else min(all_values)
    if percent_axis:
        max_value = max(max(all_values), 0.10)
    else:
        max_value = max(all_values) or 1
    if min_value == max_value:
        max_value += 1
    _add_rect(slide, plot_x, plot_y + plot_h, plot_w, 0.01, BORDER, BORDER)
    for name, values, color in series:
        clean_values = [float(value or 0) for value in values]
        if len(clean_values) < 2:
            continue
        points = []
        for index, value in enumerate(clean_values):
            px = plot_x + (plot_w * index / max(1, len(clean_values) - 1))
            py = plot_y + plot_h - (plot_h * (value - min_value) / (max_value - min_value))
            points.append((px, py))
        for start, end in zip(points, points[1:]):
            connector = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Inches(start[0]),
                Inches(start[1]),
                Inches(end[0]),
                Inches(end[1]),
            )
            connector.line.color.rgb = _rgb(color)
            connector.line.width = Pt(2)
        _add_text(slide, name, x + 0.2, y + h - 0.28 - 0.18 * series.index((name, values, color)), 1.4, 0.18, 7, color, bold=True)
    if labels:
        _add_text(slide, labels[0], plot_x, y + h - 0.2, 1.0, 0.15, 7, MUTED)
        _add_text(slide, labels[-1], plot_x + plot_w - 1.0, y + h - 0.2, 1.0, 0.15, 7, MUTED, align=PP_ALIGN.RIGHT)
    max_label = _pct(max_value) if percent_axis else _money(max_value)
    _add_text(slide, max_label, plot_x + plot_w - 1.2, plot_y - 0.08, 1.2, 0.15, 7, MUTED, align=PP_ALIGN.RIGHT)


def _add_bullets(slide, items: list[str], x: float, y: float, w: float, h: float, font_size: int = 11) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = item
        paragraph.level = 0
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = _rgb(TEXT)
        paragraph.space_after = Pt(4)


def _add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    size: int,
    color: str,
    *,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    paragraph.font.size = Pt(size)
    paragraph.font.color.rgb = _rgb(color)
    paragraph.font.bold = bold


def _add_rect(slide, x: float, y: float, w: float, h: float, fill: str, line: str) -> None:
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = _rgb(fill)
    rect.line.color.rgb = _rgb(line)
    rect.line.width = Pt(0.5)


def _share_series(dfc: list[float], total: list[float]) -> list[float]:
    return [(dfc_value / total_value if total_value else 0) for dfc_value, total_value in zip(dfc, total)]


def _money(value: float | int | None) -> str:
    return f"${float(value or 0):,.0f}"


def _pct(value: float | int | None) -> str:
    return f"{float(value or 0):.1%}"


def _signed_pct(value: float | None) -> str:
    if value is None:
        return "-"
    return f"{value:+.1%}"


def _short(value: str, limit: int) -> str:
    return value if len(value) <= limit else value[: limit - 3] + "..."


def _rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)
