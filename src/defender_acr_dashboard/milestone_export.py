from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


NAVY = "0F3A5F"
DARK = "201F1E"
TEXT = "323130"
MUTED = "605E5C"
LIGHT_BG = "F5F7FB"
BORDER = "D9E2EC"
RED = "D13438"
ORANGE = "FF8C00"
GREEN = "107C10"
BLUE = "0078D4"
WHITE = "FFFFFF"


def create_milestone_powerpoint(model: dict[str, Any], output_path: Path) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.title = "Defender milestone attach gaps"

    _title_slide(prs, model)
    _summary_slide(prs, model)
    _top_gaps_slide(prs, model)
    _methodology_slide(prs, model)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return output_path


def _title_slide(prs: Presentation, model: dict[str, Any]) -> None:
    slide = _blank_slide(prs, NAVY)
    summary = model["summary"]
    sources = model.get("sources", {})
    _add_text(slide, "Defender milestone attach gaps", 0.6, 0.55, 8.8, 0.65, 32, WHITE, bold=True)
    _add_text(slide, "Migration milestones compared with Defender for Cloud milestones", 0.62, 1.24, 7.5, 0.35, 15, "C7E0F4")
    _add_text(
        slide,
        f"Reference date: {model.get('reference_date', '-')} | Near-term window: {model.get('near_term_days', '-')} days",
        0.62,
        6.55,
        7.8,
        0.3,
        10,
        "C7E0F4",
    )
    _add_text(
        slide,
        f"Sources: {sources.get('migration', '-')} + {sources.get('defender', '-')}",
        0.62,
        6.85,
        10.8,
        0.25,
        9,
        "C7E0F4",
    )
    _add_stat_card(slide, "Accounts with gaps", f"{summary['total_accounts_with_gaps']:,}", 0.75, 2.25, 3.4, 1.45, RED)
    _add_stat_card(slide, "Opportunities with gaps", f"{summary['total_opportunities_with_gaps']:,}", 4.95, 2.25, 3.4, 1.45, ORANGE)
    _add_stat_card(slide, "Attached accounts", f"{summary['attached_accounts']:,}", 9.15, 2.25, 3.4, 1.45, GREEN)
    _add_text(
        slide,
        "Opportunity-level gaps use a strict same-account and same-Opportunity-ID comparison.",
        0.75,
        4.45,
        11.8,
        0.5,
        20,
        WHITE,
        bold=True,
    )


def _summary_slide(prs: Presentation, model: dict[str, Any]) -> None:
    slide = _blank_slide(prs)
    _slide_title(slide, "Gap summary")
    summary = model["summary"]
    priority_counts = model["priority_counts"]
    gap_type_counts = model["gap_type_counts"]

    _add_stat_card(slide, "Migration accounts", f"{summary['migration_accounts']:,}", 0.55, 1.1, 2.6, 0.95, BLUE)
    _add_stat_card(slide, "Defender accounts", f"{summary['defender_accounts']:,}", 3.35, 1.1, 2.6, 0.95, GREEN)
    _add_stat_card(slide, "Account gaps", f"{summary['account_level_gap_accounts']:,}", 6.15, 1.1, 2.6, 0.95, RED)
    _add_stat_card(slide, "Opportunity gaps", f"{summary['opportunity_level_gaps']:,}", 8.95, 1.1, 2.6, 0.95, ORANGE)

    _add_text(slide, "Priority mix", 0.75, 2.65, 3.0, 0.3, 16, DARK, bold=True)
    _add_bar_list(
        slide,
        [
            ("HIGH", priority_counts.get("HIGH", 0), RED),
            ("MEDIUM", priority_counts.get("MEDIUM", 0), ORANGE),
            ("LOW", priority_counts.get("LOW", 0), GREEN),
        ],
        0.75,
        3.05,
        5.4,
        1.45,
    )

    _add_text(slide, "Gap type mix", 6.8, 2.65, 3.0, 0.3, 16, DARK, bold=True)
    _add_bar_list(
        slide,
        [
            ("Account-level gap", gap_type_counts.get("Account-level gap", 0), RED),
            ("Opportunity-level gap", gap_type_counts.get("Opportunity-level gap", 0), ORANGE),
        ],
        6.8,
        3.05,
        5.4,
        1.1,
    )

    _add_text(
        slide,
        "High priority means at least one committed migration milestone or an estimated date inside the near-term window.",
        0.75,
        6.55,
        11.4,
        0.3,
        10,
        MUTED,
    )


def _top_gaps_slide(prs: Presentation, model: dict[str, Any]) -> None:
    slide = _blank_slide(prs)
    _slide_title(slide, "Top 10 highest priority gaps")
    headers = ["#", "Account", "Opportunity", "Gap type", "Workload", "Date", "Priority"]
    rows = [
        [
            str(index),
            _short(row["account"], 26),
            _short(row["opportunity_id"], 14),
            "Account" if row["gap_type"].startswith("Account") else "Opportunity",
            _short(row["workload"], 34),
            row["estimated_date"] or "-",
            row["priority"],
        ]
        for index, row in enumerate(model.get("top_gaps", []), start=1)
    ]
    _add_table(
        slide,
        headers,
        rows,
        0.35,
        1.15,
        12.65,
        5.8,
        [0.35, 2.35, 1.3, 1.15, 4.25, 1.0, 0.8],
    )


def _methodology_slide(prs: Presentation, model: dict[str, Any]) -> None:
    slide = _blank_slide(prs)
    _slide_title(slide, "Methodology and interpretation")
    bullets = [
        "Account-level gap: an account has Migration milestones but no Defender for Cloud milestones in the Defender workbook.",
        "Attached account: an account appears in both Migration and Defender milestone workbooks.",
        "Opportunity-level gap: for attached accounts, a Migration Opportunity ID has no Defender milestone with the same Opportunity ID.",
        f"HIGH priority: committed migration milestone or estimated date within {model.get('near_term_days', '-')} days of the reference date.",
        "MEDIUM priority: uncommitted migration milestone with a recognized workload.",
        "LOW priority: unclear or edge-case workload.",
        "Strict Opportunity ID matching may overstate gaps if Migration and Defender work are tracked under separate CRM opportunities.",
    ]
    _add_bullets(slide, bullets, 0.75, 1.2, 11.8, 4.8)


def _blank_slide(prs: Presentation, fill: str = WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = _rgb(fill)
    return slide


def _slide_title(slide, title: str) -> None:
    _add_text(slide, title, 0.55, 0.45, 8.0, 0.45, 24, DARK, bold=True)


def _add_stat_card(slide, label: str, value: str, x: float, y: float, w: float, h: float, accent: str) -> None:
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(WHITE)
    shape.line.color.rgb = _rgb(BORDER)
    bar = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(0.08), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = _rgb(accent)
    bar.line.fill.background()
    _add_text(slide, label.upper(), x + 0.18, y + 0.16, w - 0.3, 0.24, 9, MUTED, bold=True)
    _add_text(slide, value, x + 0.18, y + 0.48, w - 0.3, 0.52, 24, DARK, bold=True)


def _add_bar_list(slide, rows: list[tuple[str, int, str]], x: float, y: float, w: float, h: float) -> None:
    max_value = max(max((value for _, value, _ in rows), default=1), 1)
    row_h = h / max(len(rows), 1)
    for index, (label, value, color) in enumerate(rows):
        top = y + index * row_h
        _add_text(slide, label, x, top, 1.65, 0.24, 10, TEXT, bold=True)
        _add_text(slide, f"{value:,}", x + w - 0.65, top, 0.6, 0.24, 10, TEXT, align=PP_ALIGN.RIGHT)
        bg = slide.shapes.add_shape(1, Inches(x + 1.8), Inches(top + 0.04), Inches(w - 2.6), Inches(0.16))
        bg.fill.solid()
        bg.fill.fore_color.rgb = _rgb(LIGHT_BG)
        bg.line.fill.background()
        bar_w = (w - 2.6) * (value / max_value)
        bar = slide.shapes.add_shape(1, Inches(x + 1.8), Inches(top + 0.04), Inches(max(bar_w, 0.02)), Inches(0.16))
        bar.fill.solid()
        bar.fill.fore_color.rgb = _rgb(color)
        bar.line.fill.background()


def _add_table(slide, headers: list[str], rows: list[list[str]], x: float, y: float, w: float, h: float, col_widths: list[float]) -> None:
    row_count = max(len(rows) + 1, 2)
    table = slide.shapes.add_table(row_count, len(headers), Inches(x), Inches(y), Inches(w), Inches(h)).table
    for col_idx, width in enumerate(col_widths):
        table.columns[col_idx].width = Inches(width)
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        _format_cell(cell, bold=True, color=WHITE, fill=NAVY, size=8)
    for row_idx, row in enumerate(rows, start=1):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(value)
            _format_cell(cell, size=7)


def _add_bullets(slide, bullets: list[str], x: float, y: float, w: float, h: float) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    text_frame = box.text_frame
    text_frame.clear()
    for index, bullet in enumerate(bullets):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.size = Pt(15)
        paragraph.font.name = "Segoe UI"
        paragraph.font.color.rgb = _rgb(TEXT)


def _add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    size: int,
    color: str,
    *,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    paragraph = box.text_frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.name = "Segoe UI"
    paragraph.font.color.rgb = _rgb(color)


def _format_cell(cell, *, bold: bool = False, color: str = TEXT, fill: str = WHITE, size: int = 8) -> None:
    cell.fill.solid()
    cell.fill.fore_color.rgb = _rgb(fill)
    cell.margin_left = Inches(0.04)
    cell.margin_right = Inches(0.04)
    cell.margin_top = Inches(0.03)
    cell.margin_bottom = Inches(0.03)
    paragraph = cell.text_frame.paragraphs[0]
    paragraph.font.bold = bold
    paragraph.font.size = Pt(size)
    paragraph.font.name = "Segoe UI"
    paragraph.font.color.rgb = _rgb(color)


def _short(value: object, limit: int) -> str:
    text = str(value or "")
    return text if len(text) <= limit else text[: max(0, limit - 1)] + "..."


def _rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)
